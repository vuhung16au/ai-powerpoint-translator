import os
import glob
from pptx import Presentation
from openai import OpenAI
from dotenv import load_dotenv
import logging
import time
from datetime import datetime
import httpx
import sys
import re
from pptx.enum.text import PP_ALIGN
from pptx.util import Pt
import csv 
import argparse
import json

# Version information
VERSION = "1.0.0"

LOG_DIR = "logs"

# Load translation prompt
TRANSLATION_PROMPT = "Translation-Prompt.md"

# Set the font name for the translation target
# TODO for Japanese only
TRANSLATION_TARGET_FONTNAME = {
    "Japanese": "Meiryo UI"
    # Add more language-specific fonts as needed
}

# Supported languages from Gemini's documentation
# https://gemini.google.com/faq?hl=en-AU
# Last verified: April 2025
SUPPORTED_LANGUAGES = [
    "Arabic", "Bengali", "Bulgarian", "Chinese", "Croatian", "Czech", "Danish", "Dutch", 
    "English", "Estonian", "Finnish", "French", "German", "Greek", "Hebrew", "Hindi", 
    "Hungarian", "Indonesian", "Italian", "Japanese", "Korean", "Latvian", "Lithuanian", 
    "Malay", "Norwegian", "Persian", "Polish", "Portuguese", "Romanian", "Russian", 
    "Serbian", "Slovak", "Slovenian", "Spanish", "Swahili", "Swedish", "Thai", "Turkish", 
    "Ukrainian", "Urdu", "Vietnamese"
]

# Load language codes from JSON file
try:
    with open(os.path.join(os.path.dirname(__file__), 'lang-code.json'), 'r', encoding='utf-8') as f:
        LANGUAGE_CODES = json.load(f)
    logging.info(f"Loaded {len(LANGUAGE_CODES)} language codes from lang-code.json")
except (FileNotFoundError, json.JSONDecodeError) as e:
    logging.warning(f"Could not load language codes from JSON file: {str(e)}")
    # Fallback to hardcoded language codes
    LANGUAGE_CODES = {
        "Chinese": "zh",
        "English": "en",
        "Vietnamese": "vi",
        "Japanese": "ja"
    }
    logging.info("Using fallback hardcoded language codes")

# Supported models with their configurations
SUPPORTED_MODELS = {
    "gemini-2.0-flash-lite": {"base_url": "https://generativelanguage.googleapis.com/v1beta/openai/", "env_key": "GEMINI_API_KEY"},
    "gpt-3.5-turbo": {"base_url": "https://api.openai.com/v1", "env_key": "OPENAI_API_KEY"},
    "gpt-4": {"base_url": "https://api.openai.com/v1", "env_key": "OPENAI_API_KEY"},
    "deepseek-chat": {"base_url": "https://api.deepseek.com/v1", "env_key": "DEEPSEEK_API_KEY"}
}

# Add constant for additional requirements file
ADDITIONAL_REQUIREMENTS = "Translation-Prompt-Additional-Requirements.md"

def parse_args():
    """Parse command line arguments."""
    parser = argparse.ArgumentParser(description="Translate PowerPoint presentations using AI.")
    
    # Basic arguments
    parser.add_argument('--tone', '-t', type=str, default='formal', 
                        help='Specify the desired tone for translation (e.g., formal, informal, technical, paper)')
    parser.add_argument('--version', '-v', action='store_true', 
                        help='Show version information and exit')
    parser.add_argument('--api_key', '-k', type=str, 
                        help='Specify the API key for the translation service')
    
    # Language arguments with multiple aliases
    parser.add_argument('--source_language', '--source-language', '-sl', type=str, default='Vietnamese', 
                        help='Specify the source language of the text')
    parser.add_argument('--target_language', '--target-language', '-tl', type=str, default='Japanese', 
                        help='Specify the target language for translation')
    parser.add_argument('--supported-languages', '-sll', action='store_true', 
                        help='List supported languages and exit')
    
    # Directory arguments
    parser.add_argument('--input_dir', '-i', type=str, default='input', 
                        help='Specify the input directory containing PowerPoint files')
    parser.add_argument('--output_dir', '-o', type=str, default='output', 
                        help='Specify the output directory for translated files')
    
    # LLM model arguments
    parser.add_argument('--llm_model', '-m', type=str, default='gemini-2.0-flash-lite', 
                        help='Specify the LLM model to use for translation')
    
    # RAG arguments
    parser.add_argument('--rag', '-rag', action='store_true', 
                        help='Enable RAG for improved translation quality')
    parser.add_argument('--rag-file', '-rf', type=str, default='rag_context.txt', 
                        help='Specify the file containing RAG context or examples')
    
    args = parser.parse_args()
    
    # Handle version display
    if args.version:
        print(f"AI PowerPoint Translator v{VERSION}")
        sys.exit(0)
        
    # Handle supported languages display
    if args.supported_languages:
        print("Supported Languages:")
        for lang in SUPPORTED_LANGUAGES:
            print(f"- {lang}")
        sys.exit(0)
        
    # Validate languages
    if args.source_language not in SUPPORTED_LANGUAGES:
        print(f"Error: Source language '{args.source_language}' not supported.")
        print(f"Supported languages: {', '.join(SUPPORTED_LANGUAGES)}")
        sys.exit(1)
        
    if args.target_language not in SUPPORTED_LANGUAGES:
        print(f"Error: Target language '{args.target_language}' not supported.")
        print(f"Supported languages: {', '.join(SUPPORTED_LANGUAGES)}")
        sys.exit(1)
        
    # Validate model
    if args.llm_model not in SUPPORTED_MODELS:
        print(f"Error: Model '{args.llm_model}' not supported.")
        print(f"Supported models: {', '.join(SUPPORTED_MODELS.keys())}")
        sys.exit(1)
        
    return args

# Set up logging first, before any other imports
def setup_logging():
    # Create logs directory if it doesn't exist
    log_dir = LOG_DIR
    if not os.path.exists(log_dir):
        os.makedirs(log_dir, exist_ok=True)
    
    # Create a timestamp for the log file
    timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
    log_file = os.path.join(log_dir, f'translation_{timestamp}.log')
    
    # Configure logging to only file
    logging.basicConfig(
        level=logging.INFO,
        format='%(asctime)s - %(levelname)s - %(message)s',
        handlers=[
            logging.FileHandler(log_file, encoding='utf-8')
        ]
    )
    return log_file

# Initialize logging first
log_file = setup_logging()
logging.info("Logging system initialized")

# Load environment variables
load_dotenv()

def initialize_client(args):
    """Initialize the AI client based on model selection."""
    model_info = SUPPORTED_MODELS.get(args.llm_model)
    
    # Use provided API key or get from environment
    api_key = args.api_key or os.getenv(model_info['env_key'])
    if not api_key:
        logging.error(f"API key not found for model {args.llm_model}")
        print(f"Error: API key not found for model {args.llm_model}. Please provide with --api_key or set {model_info['env_key']} environment variable.")
        sys.exit(1)
    
    # Create a custom HTTP client
    custom_http_client = httpx.Client()
    
    # Initialize and return the client
    return OpenAI(
        api_key=api_key,
        base_url=model_info['base_url'],
        http_client=custom_http_client
    )

def load_rag_context(rag_file):
    """Load RAG context from file."""
    if not os.path.exists(rag_file):
        logging.warning(f"RAG context file not found: {rag_file}")
        return ""
    
    try:
        with open(rag_file, 'r', encoding='utf-8') as f:
            return f.read().strip()
    except Exception as e:
        logging.error(f"Error loading RAG context: {str(e)}")
        return ""

def load_additional_requirements(args):
    """Load and format additional requirements template with user arguments."""
    try:
        with open(ADDITIONAL_REQUIREMENTS, 'r', encoding='utf-8') as f:
            template = f.read()
            
        # Format the template with user arguments
        formatted = template.format(**{
            'source-language': args.source_language,
            'target-language': args.target_language,
            'tone': args.tone
        })
        return formatted
    except FileNotFoundError:
        logging.warning(f"Additional requirements file '{ADDITIONAL_REQUIREMENTS}' not found. Continuing without it.")
        return ""
    except Exception as e:
        logging.warning(f"Error loading additional requirements: {str(e)}. Continuing without it.")
        return ""

def batch_texts(texts, batch_size=30):
    """Group texts into batches for translation."""
    return [texts[i:i + batch_size] for i in range(0, len(texts), batch_size)]

def translate_batch(texts, args, client, rag_context="", max_retries=3, retry_delay=5):
    """Translate a batch of texts with the specified language, tone, and model."""
    if not texts:
        return []
    
    # Prepare the translation prompt with specified languages and tone
    with open(TRANSLATION_PROMPT, 'r', encoding='utf-8') as f:
        prompt_template = f.read()
    
    # Load additional requirements
    additional_requirements = load_additional_requirements(args)
    
    # Add RAG context if enabled
    rag_instruction = ""
    if args.rag and rag_context:
        rag_instruction = f"\nUse the following context as reference for terminology and style:\n{rag_context}\n"
    
    # Format the main prompt
    prompt = prompt_template.format(
        source_language=args.source_language,
        target_language=args.target_language,
        tone=args.tone,
        texts="\n---\n".join(texts)
    )
    
    # Combine all prompt components
    full_prompt = f"{prompt}\n\n{additional_requirements}\n{rag_instruction}".strip()
    
    # Log the request
    logging.info("=== Translation Request ===")
    logging.info(f"Source language: {args.source_language}")
    logging.info(f"Target language: {args.target_language}")
    logging.info(f"Tone: {args.tone}")
    logging.info(f"Model: {args.llm_model}")
    logging.info(f"RAG enabled: {args.rag}")
    logging.info(f"Additional requirements added: {bool(additional_requirements)}")
    for idx, text in enumerate(texts):
        logging.info(f"Text {idx + 1}: {text}")
    logging.info("=== End Request ===\n")
    
    retries = 0
    while retries <= max_retries:
        try:
            response = client.chat.completions.create(
                model=args.llm_model,
                n=1,
                messages=[
                    {"role": "system", "content": f"You are a professional translator from {args.source_language} to {args.target_language}."},
                    {"role": "user", "content": full_prompt}
                ],
                temperature=0.3,
                stream=False
            )
            
            # Log the response
            logging.info("=== Translation Response ===")
            logging.info(f"Raw response: {response.choices[0].message.content}")
            logging.info("=== End Response ===\n")
            
            # Parse the response to get translations
            translations = response.choices[0].message.content.strip().split("\n---\n")
            
            # Log parsed translations
            logging.info("=== Parsed Translations ===")
            for idx, (text, trans) in enumerate(zip(texts, translations)):
                logging.info(f"Text {idx + 1}:")
                logging.info(f"Original: {text}")
                logging.info(f"Translated: {trans}")
                logging.info("---")
            logging.info("=== End Parsed Translations ===\n")
            
            # Ensure we have the same number of translations as input texts
            if len(translations) != len(texts):
                logging.warning(f"Received {len(translations)} translations for {len(texts)} texts")
                # Pad or trim translations to match input count
                if len(translations) < len(texts):
                    translations.extend([""] * (len(texts) - len(translations)))
                    logging.warning("Added empty translations to match input count")
                else:
                    translations = translations[:len(texts)]
                    logging.warning("Trimmed excess translations")
            
            return translations
            
        except httpx.HTTPStatusError as e:
            retries += 1
            if retries > max_retries:
                logging.error(f"HTTP error after {max_retries} retries: {e.response.status_code} - {e.response.text}")
                raise
            logging.warning(f"HTTP error: {e.response.status_code}. Retrying ({retries}/{max_retries})...")
            time.sleep(retry_delay)
            
        except httpx.RequestError as e:
            retries += 1
            if retries > max_retries:
                logging.error(f"Request error after {max_retries} retries: {str(e)}")
                raise
            logging.warning(f"Request error: {str(e)}. Retrying ({retries}/{max_retries})...")
            time.sleep(retry_delay)
            
        except Exception as e:
            logging.error(f"Unexpected error during translation: {str(e)}")
            # Include traceback for debugging
            import traceback
            logging.error(f"Traceback: {traceback.format_exc()}")
            raise

def save_presentation(prs, original_filename, output_dir, target_language="Japanese"):
    """Save presentation with error handling and unique filename."""
    # Create output directory if it doesn't exist
    os.makedirs(output_dir, exist_ok=True)
    
    # Generate base output filename
    base_name = os.path.basename(original_filename)
    name_without_ext = os.path.splitext(base_name)[0]
    
    # Get language code for the target language
    lang_code = LANGUAGE_CODES.get(target_language, "unknown")
    
    # Try to save with different names if file exists or is locked
    counter = 1
    while True:
        if counter == 1:
            output_filename = os.path.join(output_dir, f"{name_without_ext}_{lang_code}.pptx")
        else:
            output_filename = os.path.join(output_dir, f"{name_without_ext}_{lang_code}_{counter}.pptx")
        
        try:
            prs.save(output_filename)
            logging.info(f"Successfully saved presentation to {output_filename}")
            return output_filename
        except PermissionError:
            logging.warning(f"Permission denied when saving to {output_filename}. File might be open in PowerPoint.")
            logging.info("Please close the file in PowerPoint if it's open.")
            counter += 1
            if counter > 5:  # Limit number of attempts
                raise Exception(f"Failed to save presentation after {counter-1} attempts. Please ensure the file is not open in PowerPoint.")
        except Exception as e:
            logging.error(f"Error saving presentation: {str(e)}")
            raise

def shape_has_valid_table(shape):
    """Check if a shape contains a valid table that can be accessed safely."""
    try:
        # Try to access the table attribute and its rows
        if hasattr(shape, "table"):
            # Access rows to validate the table is accessible
            _ = shape.table.rows
            return True
        return False
    except Exception:
        # Any exception means no valid table
        return False

def extract_table_texts(shape):
    """Extract texts from a table shape."""
    texts = []
    locations = []
    
    if not shape_has_valid_table(shape):
        return texts, locations
        
    for row_idx, row in enumerate(shape.table.rows):
        for cell_idx, cell in enumerate(row.cells):
            if cell.text.strip():
                # Process each paragraph in the cell separately
                para_texts = split_text_by_paragraphs(cell.text)
                for i, para_text in enumerate(para_texts):
                    texts.append(para_text)
                    locations.append((row_idx, cell_idx, i))  # Include paragraph index
    
    return texts, locations

def split_text_by_paragraphs(text):
    """Split text into paragraphs, handling bullet points and line breaks."""
    # First, split by line breaks
    lines = text.split('\n')
    result = []
    current_text = ""
    
    for line in lines:
        line = line.strip()
        if not line:
            if current_text:
                result.append(current_text)
                current_text = ""
            continue
            
        # Check if line starts with bullet or numbering
        if re.match(r'^[•\-\*]|\d+[.)]', line):
            # This is likely a new bullet or numbered item
            if current_text:
                result.append(current_text)
            current_text = line
        elif current_text:
            # Check if previous line had a bullet and this is continuation
            if re.match(r'^[•\-\*]|\d+[.)]', current_text.split('\n')[0]):
                current_text += '\n' + line
            else:
                # Likely a separate paragraph
                result.append(current_text)
                current_text = line
        else:
            current_text = line
    
    # Add the last paragraph if any
    if current_text:
        result.append(current_text)
        
    return result

def log_translations_to_csv(input_file, original_texts, translated_texts):
    """Log original and translated texts to a CSV file."""
    # Create logs directory if it doesn't exist
    os.makedirs(LOG_DIR, exist_ok=True)
    
    # Generate CSV file name based on input file name and timestamp
    base_name = os.path.basename(input_file)
    name_without_ext = os.path.splitext(base_name)[0]
    timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
    csv_file = os.path.join(LOG_DIR, f"{name_without_ext}-translation-log-{timestamp}.csv")
    
    # Write to CSV with double quotes
    with open(csv_file, mode='w', encoding='utf-8', newline='') as file:
        writer = csv.writer(file, quoting=csv.QUOTE_ALL)  # Ensure all fields are quoted
        writer.writerow(["Original Text", "Translated Text"])  # Write header
        for original, translated in zip(original_texts, translated_texts):
            writer.writerow([original, translated])
    
    logging.info(f"Translation log saved to {csv_file}")

def process_presentation(input_file, args, client, rag_context=""):
    """Process a PowerPoint presentation, translating text based on user preferences."""
    logging.info(f"Processing {input_file}")
    
    try:
        # Try to open the presentation file
        try:
            prs = Presentation(input_file)
        except Exception as e:
            logging.error(f"Failed to open presentation {input_file}: {str(e)}")
            import traceback
            logging.error(f"Traceback: {traceback.format_exc()}")
            raise
            
        all_texts = []
        text_locations = []
        
        # Extract all texts that need translation
        for slide_idx, slide in enumerate(prs.slides):
            for shape_idx, shape in enumerate(slide.shapes):
                # Handle regular text shapes
                if hasattr(shape, "text") and shape.text.strip():
                    # Split text by paragraphs or bullet points
                    paragraphs = []
                    
                    # Extract actual paragraphs from the text frame
                    if hasattr(shape, "text_frame"):
                        for para_idx, para in enumerate(shape.text_frame.paragraphs):
                            if para.text.strip():
                                paragraphs.append(para.text.strip())
                                text_locations.append(("paragraph", slide_idx, shape_idx, para_idx))
                    else:
                        # Fallback: split by line breaks and potential bullets
                        split_texts = split_text_by_paragraphs(shape.text)
                        for para_text in split_texts:
                            paragraphs.append(para_text)
                            # Use a special index since we don't have actual paragraph objects
                            text_locations.append(("text", slide_idx, shape_idx, len(all_texts)))
                    
                    all_texts.extend(paragraphs)
                
                # Handle tables - improved error handling with new helper function
                if shape_has_valid_table(shape):
                    try:
                        table_texts, table_locations = extract_table_texts(shape)
                        if table_texts:  # Only process if we actually got texts
                            all_texts.extend(table_texts)
                            
                            # Convert table locations to our standard format
                            for (row_idx, cell_idx, para_idx) in table_locations:
                                text_locations.append(("table", slide_idx, shape_idx, row_idx, cell_idx, para_idx))
                    except Exception as e:
                        logging.warning(f"Skipping table at slide {slide_idx+1}, shape {shape_idx}: {str(e)}")
        
        if not all_texts:
            logging.info(f"No text found in {input_file}")
            return
        
        # Translate texts in batches
        translated_texts = []
        batches = batch_texts(all_texts)
        
        print(f"\nTranslating {os.path.basename(input_file)}:")
        print(f"From {args.source_language} to {args.target_language} with '{args.tone}' tone")
        for i, batch in enumerate(batches):
            progress = (i + 1) / len(batches) * 100
            sys.stdout.write(f"\rProgress: [{int(progress)}%] Batch {i+1}/{len(batches)}")
            sys.stdout.flush()
            
            logging.info(f"Translating batch {i+1}/{len(batches)} (size: {len(batch)} texts)")
            translations = translate_batch(batch, args, client, rag_context)
            translated_texts.extend(translations)
            
            if i < len(batches) - 1:
                time.sleep(2)
        print("\nTranslation completed!")
        
        # Log translations to CSV
        log_translations_to_csv(input_file, all_texts, translated_texts)
        
        # Update presentation with translations
        for location, translated_text in zip(text_locations, translated_texts):
            if location[0] == "paragraph":
                _, slide_idx, shape_idx, para_idx = location
                shape = prs.slides[slide_idx].shapes[shape_idx]
                if hasattr(shape, "text_frame") and para_idx < len(shape.text_frame.paragraphs):
                    paragraph = shape.text_frame.paragraphs[para_idx]
                    
                    # Store original formatting
                    original_alignment = paragraph.alignment
                    original_level = paragraph.level
                    has_bullet = False
                    if hasattr(paragraph, "format") and hasattr(paragraph.format, "bullet"):
                        has_bullet = True

                    # Store original font sizes before updating text
                    original_font_sizes = []
                    for run in paragraph.runs:
                        if hasattr(run, "font") and hasattr(run.font, "size"):
                            original_font_sizes.append(run.font.size)
                        else:
                            original_font_sizes.append(None)  # None means use default
                    
                    paragraph.text = translated_text
                    
                    # Only set font if target language has a specific font defined
                    target_font = TRANSLATION_TARGET_FONTNAME.get(args.target_language)
                    if target_font:
                        # Set font for all runs in the paragraph while keeping original size
                        for idx, run in enumerate(paragraph.runs):
                            run.font.name = target_font
                            # If we have stored a font size and have enough runs, use the original
                            if idx < len(original_font_sizes) and original_font_sizes[idx] is not None:
                                run.font.size = original_font_sizes[idx]
                    
                    # Restore original formatting
                    paragraph.alignment = original_alignment
                    paragraph.level = original_level
                    if has_bullet and hasattr(paragraph, "format"):
                        try:
                            paragraph.format.bullet.enable = True
                        except:
                            # If bullet restoration fails, log but continue
                            logging.warning("Failed to restore bullet formatting")
            
            elif location[0] == "text":
                # This is a fallback for text without proper paragraph objects
                _, slide_idx, shape_idx, _ = location
                # Implementation would depend on how to handle this edge case
                pass
                
            elif location[0] == "table":
                _, slide_idx, shape_idx, row_idx, cell_idx, para_idx = location
                shape = prs.slides[slide_idx].shapes[shape_idx]
                if hasattr(shape, "table"):
                    cell = shape.table.rows[row_idx].cells[cell_idx]
                    if hasattr(cell, "text_frame") and para_idx < len(cell.text_frame.paragraphs):
                        paragraph = cell.text_frame.paragraphs[para_idx]
                        
                        # Store original formatting
                        original_alignment = paragraph.alignment
                        original_level = paragraph.level
                        has_bullet = False
                        if hasattr(paragraph, "format") and hasattr(paragraph.format, "bullet"):
                            has_bullet = True
                        
                        # Store original font sizes before updating text
                        original_font_sizes = []
                        for run in paragraph.runs:
                            if hasattr(run, "font") and hasattr(run.font, "size"):
                                original_font_sizes.append(run.font.size)
                            else:
                                original_font_sizes.append(None)  # None means use default
                        
                        paragraph.text = translated_text
                        
                        # Only set font if target language has a specific font defined
                        target_font = TRANSLATION_TARGET_FONTNAME.get(args.target_language)
                        if target_font:
                            # Set font for all runs in the paragraph while keeping original size
                            for idx, run in enumerate(paragraph.runs):
                                run.font.name = target_font
                                # If we have stored a font size and have enough runs, use the original
                                if idx < len(original_font_sizes) and original_font_sizes[idx] is not None:
                                    run.font.size = original_font_sizes[idx]
                        
                        # Restore original formatting
                        paragraph.alignment = original_alignment
                        paragraph.level = original_level
                        if has_bullet and hasattr(paragraph, "format"):
                            try:
                                paragraph.format.bullet.enable = True
                            except:
                                # If bullet restoration fails, log but continue
                                logging.warning("Failed to restore bullet formatting")
        
        # Save translated presentation with error handling
        save_presentation(prs, input_file, args.output_dir, args.target_language)
        
    except Exception as e:
        logging.error(f"Error processing presentation {input_file}: {str(e)}")
        # Include traceback for debugging
        import traceback
        logging.error(f"Traceback: {traceback.format_exc()}")
        raise

def main():
    # Parse command line arguments
    args = parse_args()
    
    # Setup logging
    log_file = setup_logging()
    logging.info(f"Translation log file: {log_file}")
    logging.info(f"Arguments: {vars(args)}")
    
    # Initialize client based on selected model
    client = initialize_client(args)
    
    # Load RAG context if enabled
    rag_context = ""
    if args.rag:
        rag_context = load_rag_context(args.rag_file)
        logging.info(f"RAG enabled, loaded context from {args.rag_file}: {len(rag_context)} characters")
    
    # Find all PPTX files in the input directory
    try:
        input_files = glob.glob(f'{args.input_dir}/*.pptx')
    except Exception as e:
        logging.error(f"Error accessing input directory: {str(e)}")
        import traceback
        logging.error(f"Traceback: {traceback.format_exc()}")
        print(f"Error accessing input directory: {str(e)}")
        return
    
    if not input_files:
        logging.warning(f"No PowerPoint files found in the input directory: {args.input_dir}")
        print(f"No PowerPoint files found in the input directory: {args.input_dir}")
        return
    
    successful_files = []
    failed_files = []
    
    for input_file in input_files:
        logging.info(f"\n=== Processing file: {input_file} ===")
        try:
            process_presentation(input_file, args, client, rag_context)
            logging.info(f"Completed translation of {input_file}")
            successful_files.append(input_file)
        except Exception as e:
            logging.error(f"Failed to process {input_file}: {str(e)}")
            failed_files.append(input_file)
            print(f"\nError processing {os.path.basename(input_file)}: {str(e)}")
            continue  # Continue with next file
    
    # Summary
    logging.info("\n=== Translation Summary ===")
    logging.info(f"Total files: {len(input_files)}")
    logging.info(f"Successfully translated: {len(successful_files)}")
    logging.info(f"Failed: {len(failed_files)}")
    
    if failed_files:
        logging.info("Failed files:")
        for file in failed_files:
            logging.info(f"- {file}")
    
    # Print summary to console as well
    print("\n=== Translation Summary ===")
    print(f"Total files: {len(input_files)}")
    print(f"Successfully translated: {len(successful_files)}")
    print(f"Failed: {len(failed_files)}")
    
    if failed_files:
        print("Failed files:")
        for file in failed_files:
            print(f"- {file}")
    
    print(f"See detailed log in: {log_file}")

if __name__ == "__main__":
    main()